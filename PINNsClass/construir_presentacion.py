"""Construye la presentacion de la clase de PINNs a partir de los resultados
realmente calculados por el cuaderno.

Requiere haber ejecutado antes `oscilador_armonico_pinn_fisica_solar.ipynb`,
que produce `figuras/` y `resultados.json`.

Uso:
    python construir_presentacion.py
"""

import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

AQUI = Path(__file__).parent
FIG = AQUI / "figuras"
SALIDA = AQUI / "Presentacion_PINNs_Fisica_Solar.pptx"

# ----------------------------------------------------------------- estilo
FONDO   = RGBColor(0x0D, 0x1B, 0x2E)
PANEL   = RGBColor(0x13, 0x25, 0x40)
PANEL2  = RGBColor(0x0A, 0x16, 0x26)
BORDE   = RGBColor(0x2C, 0x47, 0x63)
TEXTO   = RGBColor(0xE6, 0xED, 0xF7)
SUAVE   = RGBColor(0x9F, 0xB3, 0xCC)
CIAN    = RGBColor(0x3F, 0xB8, 0xF5)
AMBAR   = RGBColor(0xF5, 0xB6, 0x42)
ROJO    = RGBColor(0xFF, 0x6B, 0x6B)
VERDE   = RGBColor(0x7F, 0xBF, 0x8F)

F_TIT = "Montserrat"
F_TXT = "Open Sans"
F_COD = "DejaVu Sans Mono"

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANCA = prs.slide_layouts[6]

datos = json.loads((AQUI / "resultados.json").read_text())
fis, met = datos["fisica"], datos["metricas"]
ent, inv, sis = datos["entrenamiento"], datos["inverso"], datos["sismologia"]
neg, pin = met["caja negra"], met["PINN"]


# ----------------------------------------------------------------- helpers
def slide():
    s = prs.slides.add_slide(BLANCA)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = FONDO
    return s


def caja(s, x, y, w, h, relleno=PANEL, borde=BORDE, radio=0.02):
    f = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    f.adjustments[0] = radio
    f.fill.solid(); f.fill.fore_color.rgb = relleno
    f.line.color.rgb = borde; f.line.width = Pt(1)
    f.shadow.inherit = False
    f.text_frame.text = ""
    return f


def texto(s, x, y, w, h, partes, tam=16, fuente=F_TXT, color=TEXTO,
          align=PP_ALIGN.LEFT, interlineado=1.15, anclaje=MSO_ANCHOR.TOP):
    """partes: str, o lista de parrafos; cada parrafo es str o lista de
    (texto, dict con bold/color/size/font)."""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anclaje
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = 0

    if isinstance(partes, str):
        partes = [partes]

    for i, parrafo in enumerate(partes):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = interlineado
        p.space_after = Pt(6)
        trozos = [(parrafo, {})] if isinstance(parrafo, str) else parrafo
        for txt, fmt in trozos:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(fmt.get("size", tam))
            r.font.name = fmt.get("font", fuente)
            r.font.bold = fmt.get("bold", False)
            r.font.italic = fmt.get("italic", False)
            r.font.color.rgb = fmt.get("color", color)
    return tb


def titulo(s, txt, sub=None, y=Inches(0.42)):
    texto(s, Inches(0.7), y, W - Inches(1.4), Inches(0.8), txt,
          tam=30, fuente=F_TIT, color=CIAN)
    if sub:
        texto(s, Inches(0.7), y + Inches(0.72), W - Inches(1.4), Inches(0.5), sub,
              tam=15, color=SUAVE)


def vinetas(s, x, y, w, h, items, tam=15, color_marca=AMBAR, sep=1.25):
    """items: lista de (encabezado|None, cuerpo)."""
    parrafos = []
    for cab, cuerpo in items:
        trozos = []
        if cab:
            trozos.append(("› ", {"color": color_marca, "bold": True}))
            trozos.append((cab + "  ", {"bold": True, "color": AMBAR}))
        else:
            trozos.append(("› ", {"color": color_marca, "bold": True}))
        trozos.append((cuerpo, {}))
        parrafos.append(trozos)
    return texto(s, x, y, w, h, parrafos, tam=tam, interlineado=sep)


KEYWORDS = {"import", "from", "def", "class", "return", "for", "in", "if",
            "else", "not", "and", "or", "with", "as", "lambda", "None",
            "True", "False", "super", "self"}


def alto_codigo(lineas, tam=11.5):
    "Alto en pulgadas que ocupa un bloque de codigo, con margenes."
    return Inches(0.34 + len(lineas)*tam*1.05*1.2/72)


def codigo(s, x, y, w, lineas, h=None, tam=11.5):
    h = h if h is not None else alto_codigo(lineas, tam)
    caja(s, x, y, w, h, relleno=PANEL2, borde=BORDE)
    tb = s.shapes.add_textbox(x + Inches(0.22), y + Inches(0.16),
                              w - Inches(0.44), h - Inches(0.32))
    tf = tb.text_frame; tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        p.space_after = 0
        cuerpo, comentario = linea, ""
        if "#" in linea:
            j = linea.index("#")
            cuerpo, comentario = linea[:j], linea[j:]
        for palabra in _tokenizar(cuerpo):
            r = p.add_run(); r.text = palabra
            r.font.size = Pt(tam); r.font.name = F_COD
            r.font.color.rgb = CIAN if palabra.strip() in KEYWORDS else TEXTO
        if comentario:
            r = p.add_run(); r.text = comentario
            r.font.size = Pt(tam); r.font.name = F_COD
            r.font.color.rgb = VERDE
    return tb


def _tokenizar(linea):
    fuera, actual = [], ""
    for ch in linea:
        if ch.isalnum() or ch == "_":
            actual += ch
        else:
            if actual:
                fuera.append(actual); actual = ""
            fuera.append(ch)
    if actual:
        fuera.append(actual)
    return fuera or [""]


def imagen(s, ruta, y, alto=None, ancho=None, x=None):
    im = Image.open(ruta); rel = im.height / im.width
    if alto is not None:
        w_px = Inches(alto / rel); h_px = Inches(alto)
    else:
        w_px = Inches(ancho); h_px = Inches(ancho * rel)
    x = x if x is not None else int((W - w_px) / 2)
    return s.shapes.add_picture(str(ruta), x, y, width=w_px, height=h_px)


def pie(s, txt):
    texto(s, Inches(0.7), H - Inches(0.62), W - Inches(1.4), Inches(0.4), txt,
          tam=11, color=SUAVE, align=PP_ALIGN.CENTER)


def tabla(s, x, y, w, h, filas, anchos=None, tam=13, cabecera=True,
          colores_col=None):
    n_f, n_c = len(filas), len(filas[0])
    forma = s.shapes.add_table(n_f, n_c, x, y, w, h)
    t = forma.table
    if anchos:
        total = sum(anchos)
        for i, a in enumerate(anchos):
            t.columns[i].width = int(w * a / total)
    for i, fila in enumerate(filas):
        t.rows[i].height = Inches(h.inches / n_f)
        for j, celda in enumerate(fila):
            c = t.cell(i, j)
            c.fill.solid()
            c.fill.fore_color.rgb = PANEL if (cabecera and i == 0) else \
                (PANEL2 if i % 2 else FONDO)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = c.margin_right = Inches(0.12)
            c.margin_top = c.margin_bottom = Inches(0.03)
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(celda)
            r.font.size = Pt(tam); r.font.name = F_TXT
            r.font.bold = (cabecera and i == 0)
            if cabecera and i == 0:
                r.font.color.rgb = CIAN
            elif colores_col and j in colores_col:
                r.font.color.rgb = colores_col[j]
            else:
                r.font.color.rgb = TEXTO
    return t


def n(x, d=3):
    return f"{x:.{d}f}"


# =================================================================== 1
s = slide()
caja(s, Inches(0.9), Inches(2.55), Inches(1.6), Pt(4), relleno=AMBAR, borde=AMBAR)
texto(s, Inches(0.9), Inches(2.95), Inches(11.0), Inches(2.0),
      "Physics-Informed\nNeural Networks", tam=50, fuente=F_TIT, color=TEXTO,
      interlineado=1.05)
texto(s, Inches(0.9), Inches(5.15), Inches(10.5), Inches(0.9),
      "Cuándo meter la física dentro de la función de pérdida "
      "cambia el resultado — y cuándo no", tam=19, color=SUAVE)
texto(s, Inches(0.9), Inches(6.25), W - Inches(1.8), Inches(0.5),
      "CLASE PRÁCTICA DE PROGRAMACIÓN · ASTROFÍSICA COMPUTACIONAL 2026-I",
      tam=13, fuente=F_TIT, color=CIAN)

# =================================================================== 2
s = slide()
titulo(s, "La pregunta que vamos a responder")
texto(s, Inches(0.7), Inches(1.5), W - Inches(1.4), Inches(0.8),
      [[("Toda presentación de PINNs enseña la fórmula. Casi ninguna dice "
         "cuándo conviene usarla. Hoy respondemos eso con un experimento "
         "controlado, no con opiniones.", {})]], tam=17, color=SUAVE)

caja(s, Inches(0.7), Inches(2.6), Inches(5.9), Inches(2.95))
texto(s, Inches(1.0), Inches(2.85), Inches(5.3), Inches(0.5),
      "PROBLEMA DIRECTO", tam=15, fuente=F_TIT, color=ROJO)
texto(s, Inches(1.0), Inches(3.35), Inches(5.3), Inches(2.4),
      [[("Conozco la ecuación y sus parámetros, quiero ", {}),
        ("u(t)", {"italic": True})],
       [("Una PINN es una mala idea.", {"bold": True, "color": TEXTO})],
       [("RK4 resuelve esto ", {}),
        (f"{ent['tiempo_pinn']/ent['t_rk4']:,.0f}× más rápido", {"bold": True, "color": AMBAR}),
        (" y con un error ", {}),
        (f"{pin['RMSE u [Mm]']/ent['err_rk4']:,.0f}× menor", {"bold": True, "color": AMBAR}),
        (". No uses una red neuronal para integrar una EDO que ya sabes integrar.", {})]],
      tam=15)

caja(s, Inches(6.9), Inches(2.6), Inches(5.7), Inches(2.95))
texto(s, Inches(7.2), Inches(2.85), Inches(5.1), Inches(0.5),
      "PROBLEMA INVERSO Y ASIMILACIÓN", tam=15, fuente=F_TIT, color=CIAN)
texto(s, Inches(7.2), Inches(3.35), Inches(5.1), Inches(2.4),
      [[("Tengo observaciones escasas, ruidosas y con huecos, y quiero el "
         "estado continuo y los parámetros físicos", {})],
       [("Aquí la física en la pérdida lo cambia todo.", {"bold": True, "color": TEXTO})],
       [("Convierte un ajuste sin sentido en una ", {}),
        ("medición", {"bold": True, "color": AMBAR}),
        (". Es el caso real en física solar: cadencia limitada, ruido, "
         "variables del plasma que nadie mide directamente.", {})]],
      tam=15)

# =================================================================== 3
s = slide()
titulo(s, 'El límite del "black box"',
       "Una red que solo ve datos no aprende física: aprende correlaciones")
vinetas(s, Inches(0.7), Inches(2.0), Inches(6.1), Inches(4.4), [
    ("Extrapola mal.", "Fuera de la distribución de entrenamiento no hay "
     "ninguna garantía; la red hace lo único que sabe, interpolar suavemente."),
    ("Sobreajusta el ruido.", "Si tiene capacidad suficiente pasa exactamente "
     "por cada punto, incluida la parte que es error de medición."),
    ("Viola leyes de conservación.", "Nada en la pérdida le impide crear masa, "
     "momento o energía. Y lo hace."),
], tam=15.5, sep=1.5)

caja(s, Inches(7.1), Inches(2.0), Inches(5.5), Inches(3.5), relleno=PANEL2, borde=ROJO)
texto(s, Inches(7.4), Inches(2.25), Inches(4.9), Inches(0.5),
      "LO QUE MEDIMOS HOY EN CLASE", tam=13, fuente=F_TIT, color=ROJO)
texto(s, Inches(7.4), Inches(2.8), Inches(4.9), Inches(2.5),
      [[("Con 14 puntos ruidosos, la red sin física llevó su error de datos a ",
         {}), (f"{ent['loss_datos_negra']:.0e}", {"bold": True, "color": AMBAR}),
        (": memorizó el ruido exactamente.", {})],
       [("Consecuencia física: fabricó un ", {}),
        (f"{100*neg['energia espuria / E(0)']:.0f} %", {"bold": True, "color": ROJO}),
        (" de la energía inicial del sistema de la nada.", {})],
       [("No es un problema estético. Es un lazo coronal que se amplifica solo, "
         "sin fuente.", {"italic": True, "color": SUAVE})]], tam=14.5)

# =================================================================== 4
s = slide()
titulo(s, "El paradigma de las PINNs",
       "La red sigue siendo una red; lo que cambia es qué se optimiza")
for i, (tit, cuerpo) in enumerate([
    ("Conocimiento de datos",
     "Aprende de mediciones empíricas y condiciones de contorno o iniciales. "
     "Tolera datos ruidosos, escasos o incompletos, que es la situación normal "
     "en observaciones solares."),
    ("Restricciones físicas",
     "La ecuación diferencial se evalúa sobre la propia red mediante "
     "diferenciación automática y su incumplimiento se penaliza. Eso restringe "
     "el espacio de soluciones y reduce cuántos datos hacen falta."),
]):
    x = Inches(0.7 + i*6.2)
    caja(s, x, Inches(2.2), Inches(5.9), Inches(2.75))
    texto(s, x + Inches(0.35), Inches(2.5), Inches(5.2), Inches(0.5), tit,
          tam=19, fuente=F_TIT, color=AMBAR)
    texto(s, x + Inches(0.35), Inches(3.15), Inches(5.2), Inches(2.4), cuerpo, tam=15)

pie(s, "La clave: los puntos de colocación no necesitan datos. La física es el dato que falta.")

# =================================================================== 5
s = slide()
titulo(s, "La función de pérdida híbrida")
caja(s, Inches(1.6), Inches(1.9), Inches(10.1), Inches(1.5), relleno=PANEL2)
texto(s, Inches(1.6), Inches(2.15), Inches(10.1), Inches(1.0),
      [[("L", {"size": 30, "font": F_TIT}), ("total", {"size": 17}),
        ("  =  ", {"size": 30}),
        ("L", {"size": 30, "font": F_TIT}), ("datos", {"size": 17}),
        ("  +  λ", {"size": 30, "color": AMBAR}),
        ("fís", {"size": 17, "color": AMBAR}),
        (" L", {"size": 30, "font": F_TIT, "color": AMBAR}),
        ("física", {"size": 17, "color": AMBAR}),
        ("  +  λ", {"size": 30}), ("ci", {"size": 17}),
        (" L", {"size": 30, "font": F_TIT}), ("ci", {"size": 17})]],
      align=PP_ALIGN.CENTER)

for i, (tit, cuerpo, col) in enumerate([
    ("Datos", "MSE entre la red y las observaciones. Solo se evalúa donde hay "
     "mediciones.", CIAN),
    ("Física", "MSE del residuo de la EDO en puntos de colocación repartidos por "
     "todo el dominio, también donde no hay datos.", AMBAR),
    ("Iniciales", "Ancla la solución: sin esto, la familia de soluciones de la "
     "EDO tiene infinitos miembros.", CIAN),
]):
    x = Inches(0.7 + i*4.1)
    caja(s, x, Inches(3.7), Inches(3.8), Inches(1.9))
    texto(s, x + Inches(0.28), Inches(3.9), Inches(3.3), Inches(0.4), tit,
          tam=16, fuente=F_TIT, color=col)
    texto(s, x + Inches(0.28), Inches(4.35), Inches(3.3), Inches(1.3), cuerpo, tam=13)

texto(s, Inches(0.7), Inches(5.85), W - Inches(1.4), Inches(1.0),
      [[("Detalle que decide si funciona: ", {"bold": True, "color": AMBAR}),
        ("adimensionalizar el residuo dividiéndolo por ω₀². Sin eso los tres "
         "términos tienen unidades distintas, λ depende de la escala temporal "
         "elegida y aparece la patología de gradientes.", {})]],
      tam=14.5, color=SUAVE)

# =================================================================== 6
s = slide()
titulo(s, "Qué NO es una PINN", "Tres malentendidos que hay que desactivar antes de programar")
for i, (tit, cuerpo) in enumerate([
    ("No es un integrador mejorado",
     "Resolver una EDO/EDP conocida con una PINN es más lento y menos preciso "
     "que un método clásico. Casi siempre. Sin excepciones cómodas."),
    ("No garantiza que se cumpla la física",
     "La pérdida física es una penalización blanda: castiga el incumplimiento, "
     "no lo prohíbe. Una PINN mal entrenada viola la ecuación tranquilamente."),
    ("No es magia con datos escasos",
     "Si la ecuación impuesta está equivocada, la red devuelve con toda "
     "confianza parámetros sin sentido. La física en la pérdida es una hipótesis."),
]):
    y = Inches(1.95 + i*1.55)
    caja(s, Inches(0.7), y, Inches(11.9), Inches(1.35), relleno=PANEL)
    texto(s, Inches(1.05), y + Inches(0.16), Inches(4.0), Inches(0.9), tit,
          tam=16, fuente=F_TIT, color=ROJO, anclaje=MSO_ANCHOR.MIDDLE)
    texto(s, Inches(5.2), y + Inches(0.16), Inches(7.1), Inches(1.0), cuerpo,
          tam=14, anclaje=MSO_ANCHOR.MIDDLE)

# =================================================================== 7
s = slide()
titulo(s, "¿Cuándo sí y cuándo no?")
tabla(s, Inches(0.7), Inches(1.65), Inches(11.9), Inches(5.0), [
    ["Situación", "¿PINN?", "Por qué"],
    ["Resolver una EDO/EDP con parámetros conocidos", "NO",
     "RK4, diferencias finitas o espectral son órdenes de magnitud mejores"],
    ["Datos escasos y ruidosos + ecuación conocida", "SÍ",
     "La física regulariza e interpola donde no hay observaciones"],
    ["Estimar parámetros físicos desde observaciones", "SÍ",
     "El caso estrella: los parámetros son variables entrenables más"],
    ["Reconstruir variables no observadas", "SÍ",
     "La ecuación acopla lo observado con lo oculto"],
    ["Geometrías complicadas, mallado difícil", "SÍ",
     "Sin malla: solo hacen falta puntos de colocación"],
    ["Dinámica rígida (stiff) o multiescala", "CON CUIDADO",
     "Patologías de gradiente conocidas; suele fallar sin trucos"],
    ["Se necesita precisión de máquina", "NO",
     "Difícil bajar de 10⁻⁵–10⁻⁶ de error relativo"],
], anchos=[4.2, 1.8, 6.0], tam=12.5,
   colores_col={1: AMBAR})

# =================================================================== 8
s = slide()
titulo(s, "Caso práctico: oscilación de un lazo coronal",
       "El oscilador amortiguado no es un juguete arbitrario: es el modo kink")
caja(s, Inches(0.7), Inches(2.35), Inches(6.4), Inches(1.15), relleno=PANEL2)
texto(s, Inches(0.7), Inches(2.6), Inches(6.4), Inches(0.7),
      [[("ü  +  2β u̇  +  ω₀² u  =  0", {"size": 26, "font": F_TIT, "color": CIAN})]],
      align=PP_ALIGN.CENTER)
vinetas(s, Inches(0.7), Inches(3.75), Inches(6.4), Inches(3.0), [
    ("u(t)", "desplazamiento transversal del eje del lazo [Mm]"),
    ("ω₀", "tensión magnética como fuerza restauradora → el período P"),
    ("β", "amortiguamiento, dominado por absorción resonante → el tiempo τ"),
], tam=14.5, sep=1.35)

caja(s, Inches(7.4), Inches(2.35), Inches(5.2), Inches(4.1), relleno=PANEL2, borde=AMBAR)
texto(s, Inches(7.7), Inches(2.6), Inches(4.6), Inches(0.4),
      "POR QUÉ IMPORTA", tam=13, fuente=F_TIT, color=AMBAR)
texto(s, Inches(7.7), Inches(3.1), Inches(4.6), Inches(3.1),
      [[("Una fulguración cercana sacude el lazo y este oscila en su modo kink.",
         {})],
       [("Se observa u(t) con TRACE o SDO/AIA, se miden P y τ, y de ahí se "
         "infiere el campo magnético coronal y la estructura fina transversal "
         "del lazo.", {})],
       [("Ninguna de esas dos cosas se mide directamente. Eso es la "
         "sismología coronal.", {"bold": True, "color": TEXTO})]], tam=14)

# =================================================================== 9
s = slide()
titulo(s, "La ley que un modelo debe respetar",
       "No juzgamos por cómo se ve la curva; juzgamos por los observables físicos")
caja(s, Inches(0.7), Inches(2.15), Inches(5.7), Inches(1.2), relleno=PANEL2)
texto(s, Inches(0.7), Inches(2.42), Inches(5.7), Inches(0.7),
      [[("E  =  ½ v²  +  ½ ω₀² u²", {"size": 22, "font": F_TIT, "color": TEXTO})]],
      align=PP_ALIGN.CENTER)
caja(s, Inches(6.9), Inches(2.15), Inches(5.7), Inches(1.2), relleno=PANEL2, borde=AMBAR)
texto(s, Inches(6.9), Inches(2.42), Inches(5.7), Inches(0.7),
      [[("Ė  =  −2β v²  ≤  0", {"size": 22, "font": F_TIT, "color": AMBAR})]],
      align=PP_ALIGN.CENTER)

texto(s, Inches(0.7), Inches(3.6), W - Inches(1.4), Inches(0.6),
      [[("Sustituyendo la ecuación de movimiento en dE/dt, los términos de "
         "fuerza restauradora se cancelan y queda solo la disipación. La ley "
         "dice dos cosas, y hay que verificar las dos:", {})]], tam=15, color=SUAVE)

for i, (tit, cuerpo) in enumerate([
    ("Signo", "La energía nunca puede aumentar. Ė > 0 en cualquier instante es "
     "físicamente imposible sin una fuente."),
    ("Magnitud", "No basta con que baje: debe bajar exactamente a la tasa 2βv², "
     "que es la potencia disipada."),
]):
    x = Inches(0.7 + i*6.2)
    caja(s, x, Inches(4.35), Inches(5.9), Inches(1.35))
    texto(s, x + Inches(0.3), Inches(4.55), Inches(5.3), Inches(0.35), tit,
          tam=15, fuente=F_TIT, color=CIAN)
    texto(s, x + Inches(0.3), Inches(4.95), Inches(5.3), Inches(0.8), cuerpo, tam=13.5)

caja(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(0.95), relleno=PANEL2, borde=AMBAR)
texto(s, Inches(1.05), Inches(6.12), Inches(11.2), Inches(0.7),
      [[("Energía espuria: ", {"bold": True, "color": AMBAR}),
        ("la fracción de E(0) que un modelo crea de la nada, "
         "∫ max(Ė + 2βv², 0) dt / E(0). Para la solución verdadera vale cero. "
         "Es nuestro detector de violaciones.", {})]], tam=14)

# =================================================================== 10
s = slide()
titulo(s, "El escenario observacional",
       f"{fis['n_datos']} puntos, σ = {fis['sigma']} Mm, y un hueco de cobertura de "
       f"{fis['gap'][1]-fis['gap'][0]:.0f} minutos")
imagen(s, FIG / "01_escenario_observacional.png", Inches(1.85), alto=4.3)
pie(s, "La pregunta interesante no es qué hace cada modelo donde hay datos, "
       "sino qué inventa donde no los hay.")

# =================================================================== 11
s = slide()
titulo(s, "Paso 1: la arquitectura", "Un MLP que actúa como ansatz continuo y diferenciable")
codigo(s, Inches(0.7), Inches(1.95), Inches(7.3), [
    "class MLP(nn.Module):",
    "    def __init__(self, hidden=32, layers=3):",
    "        super().__init__()",
    "        bloques = [nn.Linear(1, hidden), nn.Tanh()]",
    "        for _ in range(layers - 1):",
    "            bloques += [nn.Linear(hidden, hidden), nn.Tanh()]",
    "        bloques += [nn.Linear(hidden, 1)]",
    "        self.net = nn.Sequential(*bloques)",
    "",
    "    def forward(self, t):",
    "        # normaliza t -> [-1,1] DENTRO del modelo:",
    "        # autograd deriva respecto al t fisico",
    "        return self.net(2.0*t/T_END - 1.0)",
])
caja(s, Inches(8.4), Inches(1.95), Inches(4.2), Inches(3.15), relleno=PANEL2, borde=AMBAR)
texto(s, Inches(8.7), Inches(2.2), Inches(3.6), Inches(0.4),
      "POR QUÉ tanh Y NO ReLU", tam=13, fuente=F_TIT, color=AMBAR)
texto(s, Inches(8.7), Inches(2.75), Inches(3.6), Inches(2.3),
      [[("Necesitamos que ü exista y sea continua.", {})],
       [("ReLU tiene segunda derivada nula en casi todo punto: el residuo de "
         "una EDO de segundo orden sería idénticamente cero salvo en un "
         "conjunto de medida nula.", {})],
       [("La pérdida física no informaría nada.", {"bold": True, "color": TEXTO})]],
      tam=13.5)

# =================================================================== 12
s = slide()
titulo(s, "Paso 2: el residuo con autograd",
       "Derivadas exactas hasta precisión de máquina, sin diferencias finitas")
codigo(s, Inches(0.7), Inches(1.95), Inches(11.9), [
    "def deriv(y, x):",
    "    # create_graph=True: hace falta para poder derivar OTRA vez",
    "    return torch.autograd.grad(y, x, torch.ones_like(y), create_graph=True)[0]",
    "",
    "def perdidas(model, lam_fis, lam_ci=20.0):",
    "    l_dat = torch.mean((model(t_data_t) - u_data_t)**2)",
    "",
    "    u_c = model(t_col_t)                    # puntos de colocacion: SIN datos",
    "    du  = deriv(u_c, t_col_t)               # velocidad",
    "    d2u = deriv(du,  t_col_t)               # aceleracion",
    "    r = (d2u + 2*beta*du + omega0**2*u_c)/omega0**2    # residuo adimensional",
    "    l_fis = torch.mean(r**2)",
    "",
    "    return l_dat + lam_fis*l_fis + lam_ci*l_ci, l_dat, l_fis, l_ci",
])

# =================================================================== 13
s = slide()
titulo(s, "Paso 3: el entrenamiento",
       "Adam explora, L-BFGS remata — el detalle que casi ningún tutorial menciona")
codigo(s, Inches(0.7), Inches(1.95), Inches(7.3), [
    "# fase 1: Adam con cosine annealing",
    "opt = torch.optim.Adam(model.parameters(), lr=5e-3)",
    "sch = CosineAnnealingLR(opt, epochs_adam, eta_min=1e-4)",
    "for ep in range(epochs_adam):",
    "    opt.zero_grad()",
    "    loss, *_ = perdidas(model, lam_fis)",
    "    loss.backward(); opt.step(); sch.step()",
    "",
    "# fase 2: L-BFGS (cuasi-Newton, usa curvatura)",
    "lbfgs = torch.optim.LBFGS(model.parameters(),",
    "            line_search_fn='strong_wolfe')",
    "def closure():",
    "    lbfgs.zero_grad()",
    "    loss, *_ = perdidas(model, lam_fis)",
    "    loss.backward()",
    "    return loss",
    "lbfgs.step(closure)",
])
caja(s, Inches(8.4), Inches(1.95), Inches(4.2), Inches(3.95), relleno=PANEL2, borde=CIAN)
texto(s, Inches(8.7), Inches(2.2), Inches(3.6), Inches(0.4),
      "EL EXPERIMENTO CONTROLADO", tam=13, fuente=F_TIT, color=CIAN)
texto(s, Inches(8.7), Inches(2.75), Inches(3.6), Inches(3.4),
      [[("Entrenamos DOS modelos con la misma red, el mismo optimizador y los "
         "mismos datos.", {})],
       [("Lo único que cambia:", {"color": SUAVE})],
       [("λfís = 0", {"bold": True, "color": ROJO}), ("  → caja negra", {})],
       [(f"λfís = {ent['lambda_fis']:.0f}", {"bold": True, "color": CIAN}),
        ("  → PINN", {})],
       [("Cualquier diferencia es atribuible a la pérdida, no al optimizador.",
         {"italic": True, "color": SUAVE})]], tam=13.5)

# =================================================================== 14
s = slide()
titulo(s, "¿Ajusta? Cinemática",
       "Mismo dato, misma red, mismo optimizador — solo cambia la pérdida")
imagen(s, FIG / "02a_cinematica.png", Inches(2.0), ancho=12.0, x=Inches(0.67))
pie(s, "A simple vista, en (a) las dos curvas pasan por los puntos. "
       "El problema aparece en las derivadas.")

# =================================================================== 15
s = slide()
titulo(s, "¿Respeta la física? Energía",
       "El diagnóstico que separa un ajuste de una reconstrucción física")
imagen(s, FIG / "02b_energia.png", Inches(2.0), ancho=12.0, x=Inches(0.67))
pie(s, "(e) es el veredicto: por encima de cero, el modelo está creando energía "
       "que ninguna fuente aporta.")

# =================================================================== 15
s = slide()
titulo(s, "Dos formas distintas de fallar")
caja(s, Inches(0.7), Inches(1.85), Inches(5.9), Inches(2.3), relleno=PANEL2, borde=ROJO)
texto(s, Inches(1.0), Inches(2.05), Inches(5.3), Inches(0.4),
      "DONDE SÍ HAY DATOS: SOBREAJUSTA EL RUIDO", tam=13, fuente=F_TIT, color=ROJO)
texto(s, Inches(1.0), Inches(2.55), Inches(5.3), Inches(1.5),
      "Para pasar exactamente por cada punto ruidoso, la curva debe curvarse "
      "mucho entre punto y punto. Curvatura es aceleración, y esas "
      "aceleraciones espurias son energía inyectada. Ajustar el ruido tiene "
      "una consecuencia física medible.", tam=14)

caja(s, Inches(6.9), Inches(1.85), Inches(5.7), Inches(2.3), relleno=PANEL2, borde=ROJO)
texto(s, Inches(7.2), Inches(2.05), Inches(5.1), Inches(0.4),
      "DONDE NO HAY DATOS: SE APLANA", tam=13, fuente=F_TIT, color=ROJO)
texto(s, Inches(7.2), Inches(2.55), Inches(5.1), Inches(1.5),
      "En el hueco la caja negra pierde la oscilación por completo: dibuja una "
      "joroba suave donde la verdad completa 1.6 ciclos. Sin datos y sin "
      "física, una red hace lo único que sabe hacer: interpolar suavemente.",
      tam=14)

tabla(s, Inches(0.7), Inches(4.35), Inches(11.9), Inches(2.5), [
    ["Diagnóstico", "Caja negra", "PINN", "Mejora"],
    ["RMSE de u  [Mm]", n(neg["RMSE u [Mm]"], 4), n(pin["RMSE u [Mm]"], 4),
     f"{neg['RMSE u [Mm]']/pin['RMSE u [Mm]']:.0f}×"],
    ["RMSE dentro del hueco  [Mm]", n(neg["RMSE u en el hueco [Mm]"], 4),
     n(pin["RMSE u en el hueco [Mm]"], 4),
     f"{neg['RMSE u en el hueco [Mm]']/pin['RMSE u en el hueco [Mm]']:.0f}×"],
    ["RMSE de v (nunca observada)", n(neg["RMSE v [Mm/min]"], 4),
     n(pin["RMSE v [Mm/min]"], 4),
     f"{neg['RMSE v [Mm/min]']/pin['RMSE v [Mm/min]']:.0f}×"],
    ["RMS del residuo de la EDO", n(neg["RMS residuo EDO"], 4),
     n(pin["RMS residuo EDO"], 4),
     f"{neg['RMS residuo EDO']/pin['RMS residuo EDO']:.0f}×"],
    ["Energía espuria / E(0)", f"{100*neg['energia espuria / E(0)']:.1f} %",
     f"{100*pin['energia espuria / E(0)']:.1f} %",
     f"{neg['energia espuria / E(0)']/pin['energia espuria / E(0)']:.0f}×"],
], anchos=[5.0, 2.3, 2.3, 1.6], tam=13,
   colores_col={1: ROJO, 2: CIAN, 3: AMBAR})

# =================================================================== 16
def gif_para_deck(origen, destino, pausa_inicial=1600, fps=8):
    """Copia el GIF anteponiendo el fotograma final, para que el PDF (y la
    vista previa) muestren el estado convergido en vez del inicial."""
    im = Image.open(origen)
    marcos = []
    for k in range(im.n_frames):
        im.seek(k)
        marcos.append(im.convert("RGBA").convert("P", palette=Image.ADAPTIVE))
    orden = [marcos[-1]] + marcos
    duraciones = [pausa_inicial] + [int(1000/fps)]*len(marcos)
    orden[0].save(destino, save_all=True, append_images=orden[1:],
                  duration=duraciones, loop=0, disposal=2)
    return im.n_frames

gif = FIG / "04_entrenamiento.gif"
gif_deck = FIG / "04_entrenamiento_deck.gif"
n_marcos = gif_para_deck(gif, gif_deck)

s = slide()
titulo(s, "Cómo la física dobla la curva durante el entrenamiento")
s.shapes.add_picture(str(gif_deck), Inches(0.75), Inches(1.75), width=Inches(11.8))
caja(s, Inches(0.75), Inches(5.6), Inches(11.8), Inches(1.15), relleno=PANEL2)
texto(s, Inches(1.05), Inches(5.78), Inches(11.2), Inches(0.85),
      [[("Ambas redes arrancan idénticas.", {"bold": True, "color": TEXTO}),
        ("  Encuentran primero los datos y luego reconstruyen la oscilación de "
         "izquierda a derecha. En el hueco se separan: la ", {}),
        ("caja negra", {"bold": True, "color": ROJO}),
        (" se queda plana porque nada la obliga a oscilar, mientras los puntos "
         "de colocación empujan a la ", {}),
        ("PINN", {"bold": True, "color": CIAN}),
        (" a continuar la oscilación a través del vacío.", {})]], tam=14)
pie(s, f"{n_marcos} instantáneas tomadas durante el entrenamiento de ambos "
       "modelos · figuras/04_entrenamiento.gif")

# =================================================================== 18
s = slide()
titulo(s, "Donde la PINN gana: el problema inverso",
       "Los parámetros físicos como variables entrenables, junto con los pesos")
imagen(s, FIG / "05_problema_inverso.png", Inches(1.75), alto=3.4)
caja(s, Inches(0.7), Inches(5.35), Inches(5.9), Inches(1.5), relleno=PANEL2, borde=CIAN)
texto(s, Inches(1.0), Inches(5.55), Inches(5.3), Inches(1.1),
      [[("ω₀ recuperada con ", {}),
        (f"{inv['err_omega0']:.1f} % de error", {"bold": True, "color": CIAN})],
       [("El período se lee de los cruces por cero: muchos ciclos lo "
         "restringen fuertemente.", {"size": 13, "color": SUAVE})]], tam=15)
caja(s, Inches(6.9), Inches(5.35), Inches(5.7), Inches(1.5), relleno=PANEL2, borde=AMBAR)
texto(s, Inches(7.2), Inches(5.55), Inches(5.1), Inches(1.1),
      [[("β recuperada con ", {}),
        (f"{inv['err_beta']:.1f} % de error", {"bold": True, "color": AMBAR})],
       [("El amortiguamiento se lee de la envolvente, mucho peor determinada "
         "con ruido y huecos. La misma jerarquía aparece en la sismología "
         "coronal real.", {"size": 13, "color": SUAVE})]], tam=15)

# =================================================================== 19
s = slide()
titulo(s, "De P y τ a la física del lazo",
       "Sismología coronal: medir lo que ningún instrumento resuelve")
caja(s, Inches(0.7), Inches(1.95), Inches(5.9), Inches(2.0), relleno=PANEL2)
texto(s, Inches(1.0), Inches(2.15), Inches(5.3), Inches(0.4),
      "CAMPO MAGNÉTICO", tam=13, fuente=F_TIT, color=CIAN)
texto(s, Inches(1.0), Inches(2.6), Inches(5.3), Inches(1.2),
      [[("c_k = 2L/P    →    B = c_k √(μ₀(ρi+ρe)/2)",
         {"font": F_COD, "size": 14, "color": TEXTO})],
       [("Velocidad kink del modo transversal en un tubo delgado.",
         {"size": 13, "color": SUAVE})]], tam=14)

caja(s, Inches(6.9), Inches(1.95), Inches(5.7), Inches(2.0), relleno=PANEL2)
texto(s, Inches(7.2), Inches(2.15), Inches(5.1), Inches(0.4),
      "ESTRUCTURA FINA TRANSVERSAL", tam=13, fuente=F_TIT, color=CIAN)
texto(s, Inches(7.2), Inches(2.6), Inches(5.1), Inches(1.2),
      [[("τ/P = (2/π)(a/l)(ζ+1)/(ζ−1)",
         {"font": F_COD, "size": 14, "color": TEXTO})],
       [("Absorción resonante en la capa inhomogénea del borde.",
         {"size": 13, "color": SUAVE})]], tam=14)

tabla(s, Inches(0.7), Inches(4.2), Inches(11.9), Inches(1.9), [
    ["Diagnóstico", "Valor verdadero", "Estimado por la PINN", "Error"],
    ["Campo magnético B", f"{sis['B_true']:.2f} G", f"{sis['B_est']:.2f} G",
     f"{sis['err_B']:.1f} %"],
    ["Ancho relativo de la capa  l/a", f"{sis['la_true']:.2f}",
     f"{sis['la_est']:.2f}", f"{sis['err_la']:.1f} %"],
], anchos=[4.4, 2.6, 2.9, 1.5], tam=14, colores_col={3: AMBAR})

texto(s, Inches(0.7), Inches(6.35), W - Inches(1.4), Inches(0.7),
      [[(f"Partiendo de {fis['n_datos']} medidas ruidosas de una posición, con un "
         f"hueco de {fis['gap'][1]-fis['gap'][0]:.0f} minutos, y sin conocer ningún "
         "parámetro del sistema.", {"bold": True})]],
      tam=15, color=TEXTO, align=PP_ALIGN.CENTER)

# =================================================================== 20
s = slide()
titulo(s, "Honestidad computacional: PINN vs RK4",
       "Las dos primeras filas son demoledoras para la PINN, y está bien que lo sean")
tabla(s, Inches(0.7), Inches(1.9), Inches(11.9), Inches(4.4), [
    ["", "RK4", "PINN"],
    ["Tiempo de cómputo", f"{1e3*ent['t_rk4']:.0f} ms", f"{ent['tiempo_pinn']:.0f} s"],
    ["Error máximo en u  [Mm]", f"{ent['err_rk4']:.1e}", f"{pin['RMSE u [Mm]']:.1e}"],
    ["¿Necesita conocer β y ω₀?", "sí", "no: los estima"],
    ["¿Usa las observaciones?", "no", "sí"],
    ["¿Tolera huecos y ruido?", "no aplica", "sí"],
    ["¿Estima parámetros desconocidos?", "no", "sí"],
    ["¿Reconstruye variables ocultas?", "no", "sí"],
], anchos=[5.5, 3.2, 3.2], tam=13.5, colores_col={1: SUAVE, 2: CIAN})
texto(s, Inches(0.7), Inches(6.45), W - Inches(1.4), Inches(0.7),
      [[("Usa una PINN cuando el problema esté mal planteado en el sentido "
         "clásico. Si está bien planteado y solo hay que integrarlo, integra.",
         {"bold": True})]], tam=15, color=AMBAR, align=PP_ALIGN.CENTER)

# =================================================================== 21
s = slide()
titulo(s, "Desafíos y fronteras abiertas", "No todo está resuelto")
vinetas(s, Inches(0.7), Inches(2.0), Inches(11.9), Inches(4.4), [
    ("Patología de gradientes.",
     "Los términos de la pérdida pueden tener magnitudes de gradiente muy "
     "distintas y la convergencia se rompe. Adimensionalizar ayuda; hay "
     "esquemas de pesos adaptativos (Wang et al. 2021)."),
    ("Costo de la retropropagación.",
     "Calcular derivadas de alto orden sobre mallas multidimensionales es "
     "intensivo en memoria de GPU. Cada orden de derivada agranda el grafo."),
    ("Dinámica multiescala y rígida.",
     "Los fenómenos astrofísicos varían enormemente en escala espacial y "
     "temporal. Las PINNs tienen modos de fallo caracterizados y reproducibles "
     "(Krishnapriyan et al. 2021): a veces la solución trivial gana."),
    ("Y el riesgo silencioso.",
     "La PINN acierta porque la ecuación impuesta era la correcta. Con física "
     "equivocada devuelve un ajuste bonito con parámetros sin sentido."),
], tam=15, sep=1.45)

# =================================================================== 22
s = slide()
titulo(s, "Qué llevarse de esta clase")
for i, (num, tit, cuerpo) in enumerate([
    ("1", "La física entra en la pérdida, no en el comentario",
     "El residuo de la EDO evaluado sobre la propia red, con autograd."),
    ("2", "Audita observables, no curvas",
     "Energía, residuo, retrato de fase. Una curva que pasa por los puntos "
     f"puede estar creando un {100*neg['energia espuria / E(0)']:.0f} % de la energía inicial."),
    ("3", "El nicho es el problema inverso",
     "Para integrar una EDO conocida, integra. Para medir parámetros desde "
     "datos incompletos, la PINN es la herramienta."),
    ("4", "Adam explora, L-BFGS remata",
     "Sin la segunda fase, la mayoría de las PINNs de tutorial se quedan a "
     "medio converger."),
]):
    y = Inches(1.85 + i*1.28)
    caja(s, Inches(0.7), y, Inches(11.9), Inches(1.1))
    texto(s, Inches(0.95), y + Inches(0.1), Inches(0.7), Inches(0.9), num,
          tam=26, fuente=F_TIT, color=AMBAR, anclaje=MSO_ANCHOR.MIDDLE)
    texto(s, Inches(1.75), y + Inches(0.12), Inches(4.6), Inches(0.9), tit,
          tam=15, fuente=F_TIT, color=CIAN, anclaje=MSO_ANCHOR.MIDDLE)
    texto(s, Inches(6.5), y + Inches(0.12), Inches(5.9), Inches(0.9), cuerpo,
          tam=13.5, anclaje=MSO_ANCHOR.MIDDLE)

# =================================================================== 23
s = slide()
titulo(s, "Referencias", y=Inches(0.6))
texto(s, Inches(0.7), Inches(1.6), Inches(5.9), Inches(0.4), "PINNs",
      tam=17, fuente=F_TIT, color=AMBAR)
texto(s, Inches(0.7), Inches(2.1), Inches(5.9), Inches(4.0),
      ["Raissi, Perdikaris & Karniadakis (2019), J. Comput. Phys. 378, 686 — el artículo fundacional.",
       "Karniadakis et al. (2021), Nature Reviews Physics 3, 422 — panorama general.",
       "Wang, Teng & Perdikaris (2021), SIAM J. Sci. Comput. — patologías de gradiente.",
       "Krishnapriyan et al. (2021), NeurIPS — modos de fallo. Lectura obligatoria."],
      tam=13.5, interlineado=1.2)
texto(s, Inches(7.1), Inches(1.6), Inches(5.5), Inches(0.4), "Sismología coronal",
      tam=17, fuente=F_TIT, color=AMBAR)
texto(s, Inches(7.1), Inches(2.1), Inches(5.5), Inches(4.0),
      ["Nakariakov et al. (1999), Science 285, 862 — oscilaciones transversales con TRACE.",
       "Nakariakov & Ofman (2001), A&A 372, L53 — campo magnético desde el período kink.",
       "Ruderman & Roberts (2002), ApJ 577, 475; Goossens et al. (2002), A&A 394, L39 — absorción resonante.",
       "Antolin & Van Doorsselaere (2019), Frontiers in Physics — estructura fina de lazos."],
      tam=13.5, interlineado=1.2)
pie(s, "Cuaderno reproducible: oscilador_armonico_pinn_fisica_solar.ipynb · "
       "todas las cifras de esta presentación salen de resultados.json")

# =================================================================== 24
s = slide()
texto(s, Inches(0.9), Inches(3.0), W - Inches(1.8), Inches(1.0), "¿Preguntas?",
      tam=48, fuente=F_TIT, color=CIAN, align=PP_ALIGN.CENTER)
texto(s, Inches(0.9), Inches(4.1), W - Inches(1.8), Inches(0.5),
      "Gracias por vuestra atención.", tam=17, color=SUAVE, align=PP_ALIGN.CENTER)
texto(s, Inches(0.9), Inches(4.7), W - Inches(1.8), Inches(0.5),
      "El cuaderno corre completo en ~3 minutos en CPU.",
      tam=14, fuente=F_TIT, color=AMBAR, align=PP_ALIGN.CENTER)


# Metadatos del archivo. python-pptx deja los suyos por defecto ("Steve Canny",
# "generated using python-pptx"); los limpiamos. Pon tu nombre en AUTOR.
AUTOR = ""

props = prs.core_properties
props.title = "Physics-Informed Neural Networks"
props.subject = "Astrofísica Computacional 2026-I — clase práctica sobre PINNs"
props.author = AUTOR
props.last_modified_by = AUTOR
props.comments = ""
props.category = ""
props.keywords = "PINNs, física solar, sismología coronal, oscilaciones kink"

prs.save(SALIDA)
print(f"Presentacion generada: {SALIDA}  ({len(prs.slides.__iter__.__self__._sldIdLst)} diapositivas)")
